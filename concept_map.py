"""Concept map layout and export utilities for the Streamlit builder."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from typing import Dict, Iterable, List, Optional, Sequence, Tuple
from datetime import datetime
import uuid

import math
import networkx as nx
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from matplotlib import patches  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR  # noqa: E402
from pptx.enum.text import MSO_AUTO_SIZE  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402
from lxml import etree  # noqa: E402


EMU_PER_IN = 914_400
EMU_PER_PT = 12_700


@dataclass
class NodeLayout:
    node_id: str
    label: str
    node_type: str
    lines: List[str]
    font_pt: int
    size_emu: Tuple[int, int]
    center_emu: Tuple[int, int]

    @property
    def width_in(self) -> float:
        return self.size_emu[0] / EMU_PER_IN

    @property
    def height_in(self) -> float:
        return self.size_emu[1] / EMU_PER_IN

    @property
    def center_in(self) -> Tuple[float, float]:
        cx, cy = self.center_emu
        return cx / EMU_PER_IN, cy / EMU_PER_IN


def wrap_to_width(text: str, max_chars_per_line: int) -> List[str]:
    if max_chars_per_line < 5:
        max_chars_per_line = 5
    words = text.split()
    lines: List[str] = []
    cur = ""
    for word in words:
        trial = (cur + " " + word).strip()
        if len(trial) <= max_chars_per_line:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            if len(word) > max_chars_per_line:
                frag = word
                while len(frag) > max_chars_per_line:
                    lines.append(frag[:max_chars_per_line])
                    frag = frag[max_chars_per_line:]
                cur = frag
            else:
                cur = word
    if cur:
        lines.append(cur)
    return lines


def rect_from_center(center_xy_emu: Tuple[int, int], size_emu: Tuple[int, int]) -> Tuple[int, int, int, int]:
    cx, cy = center_xy_emu
    w, h = size_emu
    return cx - w // 2, cy - h // 2, w, h


def rects_overlap(
    r1: Tuple[int, int, int, int],
    r2: Tuple[int, int, int, int],
    pad_emu: int = 0,
) -> bool:
    x1, y1, w1, h1 = r1
    x2, y2, w2, h2 = r2
    x1 -= pad_emu
    y1 -= pad_emu
    w1 += 2 * pad_emu
    h1 += 2 * pad_emu
    x2 -= pad_emu
    y2 -= pad_emu
    w2 += 2 * pad_emu
    h2 += 2 * pad_emu
    return not (
        x1 + w1 <= x2
        or x2 + w2 <= x1
        or y1 + h1 <= y2
        or y2 + h2 <= y1
    )


def find_non_overlapping_center(
    target_c: Tuple[int, int],
    size_emu: Tuple[int, int],
    placed_rects: List[Tuple[int, int, int, int]],
    bbox: Tuple[int, int, int, int],
    pad_emu: int,
    max_radius_emu: int,
    step_emu: int,
) -> Tuple[int, int]:
    tx, ty = target_c
    w, h = size_emu
    left, top, right, bottom = bbox
    tx = max(left + w // 2 + pad_emu, min(right - w // 2 - pad_emu, tx))
    ty = max(top + h // 2 + pad_emu, min(bottom - h // 2 - pad_emu, ty))

    trial = rect_from_center((tx, ty), size_emu)
    if not any(rects_overlap(trial, rect, pad_emu) for rect in placed_rects):
        return tx, ty

    r = step_emu
    angle_step = math.radians(15)
    while r <= max_radius_emu:
        steps = max(24, int(2 * math.pi / angle_step))
        for i in range(steps):
            ang = i * angle_step
            cx = int(tx + r * math.cos(ang))
            cy = int(ty + r * math.sin(ang))
            cx = max(left + w // 2 + pad_emu, min(right - w // 2 - pad_emu, cx))
            cy = max(top + h // 2 + pad_emu, min(bottom - h // 2 - pad_emu, cy))
            trial = rect_from_center((cx, cy), size_emu)
            if not any(rects_overlap(trial, rect, pad_emu) for rect in placed_rects):
                return cx, cy
        r += step_emu

    return tx, ty


def measure_box_for_text(
    text: str,
    font_max_pt: int,
    font_min_pt: int,
    target_width_in: float,
    min_width_in: float,
    max_width_in: float,
    max_lines: int = 3,
    char_w_factor: float = 0.52,
    line_h_factor: float = 1.12,
    pad_in: float = 0.07,
) -> Tuple[List[str], int, int, int]:
    for font in range(font_max_pt, font_min_pt - 1, -1):
        inner_w_emu = max(int((target_width_in - 2 * pad_in) * EMU_PER_IN), 1)
        chars_per_line = max(int(inner_w_emu / (font * EMU_PER_PT * char_w_factor)), 8)
        lines = wrap_to_width(text, chars_per_line)
        if len(lines) <= max_lines:
            longest = max((len(line) for line in lines), default=1)
            text_w_emu = int(longest * font * EMU_PER_PT * char_w_factor)
            width_emu = text_w_emu + int(2 * pad_in * EMU_PER_IN)
            width_emu = max(
                int(min_width_in * EMU_PER_IN),
                min(int(max_width_in * EMU_PER_IN), width_emu),
            )
            chars_per_line = max(
                int((width_emu - int(2 * pad_in * EMU_PER_IN)) / (font * EMU_PER_PT * char_w_factor)),
                8,
            )
            lines = wrap_to_width(text, chars_per_line)
            line_h = int(font * EMU_PER_PT * line_h_factor)
            height_emu = int(len(lines) * line_h + 2 * pad_in * EMU_PER_IN)
            return lines, font, width_emu, height_emu

    font = font_min_pt
    chars_per_line = max(
        int(((target_width_in - 2 * pad_in) * EMU_PER_IN) / (font * EMU_PER_PT * char_w_factor)),
        8,
    )
    lines = wrap_to_width(text, chars_per_line)
    longest = max((len(line) for line in lines), default=1)
    width_emu = int(longest * font * EMU_PER_PT * char_w_factor + 2 * pad_in * EMU_PER_IN)
    width_emu = max(
        int(min_width_in * EMU_PER_IN),
        min(int(max_width_in * EMU_PER_IN), width_emu),
    )
    line_h = int(font * EMU_PER_PT * 1.1)
    height_emu = int(len(lines) * line_h + 2 * pad_in * EMU_PER_IN)
    return lines, font, width_emu, height_emu


class ConceptMapBuilder:
    """Encapsulates layout logic and export helpers for concept maps."""

    def __init__(
        self,
        slide_size: Tuple[float, float] = (13.333, 7.5),
        margins: Tuple[float, float, float, float] = (0.3, 0.3, 0.3, 0.3),
        concept_fill: Tuple[int, int, int] = (245, 245, 245),
        concept_line: Tuple[int, int, int] = (120, 120, 120),
        text_rgb: Tuple[int, int, int] = (50, 50, 50),
    ) -> None:
        self.slide_size = slide_size
        self.margins = margins
        self.concept_fill = concept_fill
        self.concept_line = concept_line
        self.text_rgb = text_rgb

    def compute_layout(
        self,
        concepts: Sequence[Dict[str, str]],
        edges: Sequence[Dict[str, str]],
        *,
        seed: int = 7,
    ) -> Dict[str, NodeLayout]:
        nodes = {
            concept["id"]: {
                "label": concept["label"],
                "type": concept.get("type", "concept"),
            }
            for concept in concepts
        }
        if not nodes:
            return {}

        graph = nx.Graph()
        graph.add_nodes_from(nodes.keys())
        for edge in edges:
            # Handle both dict format (new) and tuple format (legacy)
            if isinstance(edge, dict):
                src = edge.get("source")
                dst = edge.get("target")
            else:
                src, dst = edge
            if src in nodes and dst in nodes:
                graph.add_edge(src, dst)

        if len(graph.nodes) == 1:
            single_id = next(iter(graph.nodes))
            pos = {single_id: (0.5, 0.5)}
        else:
            base_k = 1.2 / max(1, math.sqrt(len(graph.nodes)))
            pos1 = nx.spring_layout(graph, seed=seed, k=base_k, iterations=300)
            try:
                pos2 = nx.kamada_kawai_layout(graph, weight=None, scale=1.0)
            except nx.NetworkXError:
                pos2 = pos1
            pos = {
                node: (
                    0.5 * pos1[node][0] + 0.5 * pos2[node][0],
                    0.5 * pos1[node][1] + 0.5 * pos2[node][1],
                )
                for node in graph.nodes
            }

        slide_width_emu = int(self.slide_size[0] * EMU_PER_IN)
        slide_height_emu = int(self.slide_size[1] * EMU_PER_IN)
        margins_emu = tuple(int(Inches(value)) for value in self.margins)
        left_m, top_m, right_m, bottom_m = margins_emu

        avail_w = slide_width_emu - (left_m + right_m)
        avail_h = slide_height_emu - (top_m + bottom_m)

        for nid, info in nodes.items():
            label = info["label"]
            lines, font_pt, w_emu, h_emu = measure_box_for_text(
                label,
                font_max_pt=9,
                font_min_pt=5,
                target_width_in=2.0,
                min_width_in=1.1,
                max_width_in=2.6,
                max_lines=3,
            )
            info["lines"] = lines
            info["font"] = font_pt
            info["size_emu"] = (int(w_emu), int(h_emu))

        xs = [coords[0] for coords in pos.values()]
        ys = [coords[1] for coords in pos.values()]
        min_x, max_x = min(xs), max(xs)
        min_y, max_y = min(ys), max(ys)

        def to_slide(pt: Tuple[float, float]) -> Tuple[int, int]:
            x, y = pt
            nx_ = (x - min_x) / (max_x - min_x + 1e-9)
            ny_ = (y - min_y) / (max_y - min_y + 1e-9)
            return int(left_m + nx_ * avail_w), int(top_m + ny_ * avail_h)

        target_centers = {node: to_slide(pos[node]) for node in pos}

        order = sorted(
            nodes.keys(),
            key=lambda node: (-graph.degree(node), -len(nodes[node]["label"])),
        )

        placed_centers: Dict[str, Tuple[int, int]] = {}
        placed_rects: List[Tuple[int, int, int, int]] = []
        pad = int(0.05 * EMU_PER_IN)
        bbox = (int(left_m), int(top_m), int(left_m + avail_w), int(top_m + avail_h))
        step_emu = int(0.15 * EMU_PER_IN)
        max_radius_emu = int(2.5 * EMU_PER_IN)

        for node_id in order:
            target = target_centers.get(node_id, (int(left_m + avail_w / 2), int(top_m + avail_h / 2)))
            size = nodes[node_id]["size_emu"]
            center = find_non_overlapping_center(target, size, placed_rects, bbox, pad, max_radius_emu, step_emu)
            placed_centers[node_id] = center
            placed_rects.append(rect_from_center(center, size))

        rx_min = min(
            cx - nodes[node_id]["size_emu"][0] // 2
            for node_id, (cx, cy) in placed_centers.items()
        )
        ry_min = min(
            cy - nodes[node_id]["size_emu"][1] // 2
            for node_id, (cx, cy) in placed_centers.items()
        )
        rx_max = max(
            cx + nodes[node_id]["size_emu"][0] // 2
            for node_id, (cx, cy) in placed_centers.items()
        )
        ry_max = max(
            cy + nodes[node_id]["size_emu"][1] // 2
            for node_id, (cx, cy) in placed_centers.items()
        )

        bw = rx_max - rx_min
        bh = ry_max - ry_min
        scale = min(
            float(avail_w - 2 * pad) / bw if bw > 0 else 1.0,
            float(avail_h - 2 * pad) / bh if bh > 0 else 1.0,
        )
        if scale > 1.05:
            for node_id, (cx, cy) in placed_centers.items():
                ncx = int(left_m + pad + (cx - rx_min) * scale)
                ncy = int(top_m + pad + (cy - ry_min) * scale)
                placed_centers[node_id] = (ncx, ncy)

        layout: Dict[str, NodeLayout] = {}
        for node_id, info in nodes.items():
            layout[node_id] = NodeLayout(
                node_id=node_id,
                label=info["label"],
                node_type=info["type"],
                lines=info["lines"],
                font_pt=info["font"],
                size_emu=info["size_emu"],
                center_emu=placed_centers[node_id],
            )

        return layout

    def build_presentation(
        self,
        concepts: Sequence[Dict[str, str]],
        edges: Sequence[Dict[str, str]],
        layout: Optional[Dict[str, NodeLayout]] = None,
    ) -> Presentation:
        if layout is None:
            layout = self.compute_layout(concepts, edges)
        prs = Presentation()
        prs.slide_width = Emu(self.slide_size[0] * EMU_PER_IN)
        prs.slide_height = Emu(self.slide_size[1] * EMU_PER_IN)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        slide_shapes: Dict[str, any] = {}
        for node in layout.values():
            cx, cy = node.center_emu
            w, h = node.size_emu
            x, y = cx - w // 2, cy - h // 2
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*self.concept_fill)
            shape.line.color.rgb = RGBColor(*self.concept_line)
            shape.line.width = Emu(1270)
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.NONE
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = "\n".join(node.lines)
            run.font.name = "Calibri"
            run.font.size = Pt(node.font_pt)
            run.font.color.rgb = RGBColor(*self.text_rgb)
            slide_shapes[node.node_id] = shape

        for edge in edges:
            # Handle both dict format (new) and tuple format (legacy)
            if isinstance(edge, dict):
                src = edge.get("source")
                dst = edge.get("target")
                edge_label = edge.get("label", "")
            else:
                src, dst = edge
                edge_label = ""
            
            if src not in slide_shapes or dst not in slide_shapes:
                continue
            sa = slide_shapes[src]
            sb = slide_shapes[dst]
            ax = int(sa.left + sa.width / 2)
            ay = int(sa.top + sa.height / 2)
            bx = int(sb.left + sb.width / 2)
            by = int(sb.top + sb.height / 2)
            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax, ay, bx, by)
            try:
                connector.begin_connect(sa, 0)
                connector.end_connect(sb, 0)
            except Exception:
                pass
            connector.line.width = Emu(1000)
            connector.line.color.rgb = RGBColor(*self.concept_line)
            ln = connector.line._get_or_add_ln()
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            head_end = etree.SubElement(ln, f"{{{ns}}}headEnd")
            head_end.set("type", "triangle")
            
            # Add label text box if edge has a label
            if edge_label and edge_label.strip():
                mid_x = (ax + bx) / 2
                mid_y = (ay + by) / 2
                label_text = edge_label.strip()
                lines, font_pt, w_emu, h_emu = measure_box_for_text(
                    label_text,
                    font_max_pt=8,
                    font_min_pt=6,
                    target_width_in=1.5,
                    min_width_in=0.5,
                    max_width_in=2.0,
                    max_lines=2,
                )
                label_x = int(mid_x - w_emu / 2)
                label_y = int(mid_y - h_emu / 2)
                label_shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    label_x, label_y, w_emu, h_emu
                )
                label_shape.fill.solid()
                label_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                label_shape.line.color.rgb = RGBColor(*self.concept_line)
                label_shape.line.width = Emu(635)
                label_frame = label_shape.text_frame
                label_frame.clear()
                label_frame.word_wrap = True
                label_frame.auto_size = MSO_AUTO_SIZE.NONE
                label_para = label_frame.paragraphs[0]
                label_run = label_para.add_run()
                label_run.text = "\n".join(lines)
                label_run.font.name = "Calibri"
                label_run.font.size = Pt(font_pt)
                label_run.font.color.rgb = RGBColor(*self.text_rgb)

        return prs

    def build_pptx_bytes(
        self,
        concepts: Sequence[Dict[str, str]],
        edges: Sequence[Dict[str, str]],
        layout: Optional[Dict[str, NodeLayout]] = None,
    ) -> BytesIO:
        prs = self.build_presentation(concepts, edges, layout=layout)
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    def render_preview(
        self,
        layout: Dict[str, NodeLayout],
        edges: Sequence[Dict[str, str]],
        *,
        figsize: Tuple[float, float] = (8, 4.8),
        show_edge_labels: bool = True,
    ) -> BytesIO:
        if not layout:
            raise ValueError("Cannot render preview without nodes.")

        fig, ax = plt.subplots(figsize=figsize)
        width_in, height_in = self.slide_size
        ax.set_xlim(0, width_in)
        ax.set_ylim(height_in, 0)  # invert Y to match PPT coordinates
        ax.axis("off")

        for edge in edges:
            # Handle both dict format (new) and tuple format (legacy)
            if isinstance(edge, dict):
                src = edge.get("source")
                dst = edge.get("target")
                edge_label = edge.get("label", "")
            else:
                src, dst = edge
                edge_label = ""
            
            if src not in layout or dst not in layout:
                continue
            sx, sy = layout[src].center_in
            dx, dy = layout[dst].center_in
            
            # Draw arrow
            mid_x = (sx + dx) / 2
            mid_y = (sy + dy) / 2
            
            ax.annotate(
                "",
                xy=(dx, dy),
                xytext=(sx, sy),
                arrowprops=dict(arrowstyle="->", color="#7a7a7a", linewidth=1.8),
            )
            
            # Add label if present and labels are enabled
            if show_edge_labels and edge_label and edge_label.strip():
                ax.text(
                    mid_x, mid_y,
                    edge_label.strip(),
                    ha="center",
                    va="center",
                    fontsize=7,
                    color="#323232",
                    fontname="Calibri",
                    bbox=dict(boxstyle="round,pad=0.3", facecolor="white", edgecolor="#787878", linewidth=0.8),
                    zorder=10
                )

        for node in layout.values():
            cx, cy = node.center_in
            width = node.width_in
            height = node.height_in
            rect = patches.FancyBboxPatch(
                (cx - width / 2, cy - height / 2),
                width,
                height,
                boxstyle="round,pad=0.05",
                linewidth=1.5,
                edgecolor="#787878",
                facecolor="#f5f5f5",
            )
            ax.add_patch(rect)
            ax.text(
                cx,
                cy,
                "\n".join(node.lines),
                ha="center",
                va="center",
                fontsize=node.font_pt,
                color="#323232",
                fontname="Calibri",
            )

        buffer = BytesIO()
        fig.savefig(buffer, format="png", dpi=150, bbox_inches="tight")
        plt.close(fig)
        buffer.seek(0)
        return buffer

    def render_jpg(
        self,
        layout: Dict[str, NodeLayout],
        edges: Sequence[Dict[str, str]],
        *,
        figsize: Tuple[float, float] = (8, 4.8),
        dpi: int = 300,
    ) -> BytesIO:
        """Render the concept map as a JPG image with high resolution."""
        if not layout:
            raise ValueError("Cannot render JPG without nodes.")

        fig, ax = plt.subplots(figsize=figsize, dpi=dpi / 72.0)
        width_in, height_in = self.slide_size
        ax.set_xlim(0, width_in)
        ax.set_ylim(height_in, 0)  # invert Y to match PPT coordinates
        ax.axis("off")
        ax.set_facecolor("white")

        for edge in edges:
            # Handle both dict format (new) and tuple format (legacy)
            if isinstance(edge, dict):
                src = edge.get("source")
                dst = edge.get("target")
                edge_label = edge.get("label", "")
            else:
                src, dst = edge
                edge_label = ""
            
            if src not in layout or dst not in layout:
                continue
            sx, sy = layout[src].center_in
            dx, dy = layout[dst].center_in
            
            # Draw arrow
            mid_x = (sx + dx) / 2
            mid_y = (sy + dy) / 2
            
            ax.annotate(
                "",
                xy=(dx, dy),
                xytext=(sx, sy),
                arrowprops=dict(arrowstyle="->", color="#7a7a7a", linewidth=2.0),
            )
            
            # Add label if present
            if edge_label and edge_label.strip():
                ax.text(
                    mid_x, mid_y,
                    edge_label.strip(),
                    ha="center",
                    va="center",
                    fontsize=8,
                    color="#323232",
                    fontname="Calibri",
                    bbox=dict(boxstyle="round,pad=0.4", facecolor="white", edgecolor="#787878", linewidth=1.0),
                    zorder=10
                )

        for node in layout.values():
            cx, cy = node.center_in
            width = node.width_in
            height = node.height_in
            rect = patches.FancyBboxPatch(
                (cx - width / 2, cy - height / 2),
                width,
                height,
                boxstyle="round,pad=0.05",
                linewidth=2.0,
                edgecolor="#787878",
                facecolor="#f5f5f5",
            )
            ax.add_patch(rect)
            ax.text(
                cx,
                cy,
                "\n".join(node.lines),
                ha="center",
                va="center",
                fontsize=max(8, node.font_pt + 1),
                color="#323232",
                fontname="Calibri",
            )

        buffer = BytesIO()
        fig.savefig(buffer, format="jpeg", dpi=dpi, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buffer.seek(0)
        return buffer

    def build_cxl_xml(
        self,
        concepts: Sequence[Dict[str, str]],
        edges: Sequence[Dict[str, str]],
        layout: Optional[Dict[str, NodeLayout]] = None,
        title: str = "Concept Map",
    ) -> BytesIO:
        """Build a CXL (Concept Mapping Extensible Language) XML file from the concept map data."""
        if layout is None:
            layout = self.compute_layout(concepts, edges)
        
        # Namespaces
        CMAP_NS = "http://cmap.ihmc.us/xml/cmap/"
        DC_NS = "http://purl.org/dc/elements/1.1/"
        DCTERMS_NS = "http://purl.org/dc/terms/"
        VCARD_NS = "http://www.w3.org/2001/vcard-rdf/3.0#"
        
        # Create root element
        root = etree.Element(
            "{" + CMAP_NS + "}cmap",
            nsmap={
                None: CMAP_NS,
                "dc": DC_NS,
                "dcterms": DCTERMS_NS,
                "vcard": VCARD_NS,
            }
        )
        
        # Resource metadata
        res_meta = etree.SubElement(root, "{" + CMAP_NS + "}res-meta")
        etree.SubElement(res_meta, "{" + DC_NS + "}title").text = title
        etree.SubElement(res_meta, "{" + DC_NS + "}description")
        etree.SubElement(res_meta, "{" + DC_NS + "}format").text = "x-cmap/x-storable"
        etree.SubElement(res_meta, "{" + DC_NS + "}type")
        
        # Creator/contributor (minimal)
        creator = etree.SubElement(res_meta, "{" + DC_NS + "}creator")
        etree.SubElement(creator, "{" + VCARD_NS + "}FN").text = "Concept Map Builder"
        
        contributor = etree.SubElement(res_meta, "{" + DC_NS + "}contributor")
        etree.SubElement(contributor, "{" + VCARD_NS + "}FN").text = "Concept Map Builder"
        
        etree.SubElement(res_meta, "{" + DC_NS + "}language").text = "en"
        etree.SubElement(res_meta, "{" + DC_NS + "}publisher").text = "Concept Map Builder Tool"
        
        now = datetime.now().isoformat()
        etree.SubElement(res_meta, "{" + DCTERMS_NS + "}modified").text = now
        etree.SubElement(res_meta, "{" + DCTERMS_NS + "}created").text = now
        
        # Calculate map dimensions
        if layout:
            # Find bounding box
            x_coords = []
            y_coords = []
            for node in layout.values():
                x, y = node.center_emu
                w, h = node.size_emu
                x_coords.extend([x - w//2, x + w//2])
                y_coords.extend([y - h//2, y + h//2])
            if x_coords and y_coords:
                map_width = max(x_coords) - min(x_coords) + 200  # padding
                map_height = max(y_coords) - min(y_coords) + 200
            else:
                map_width = 2200
                map_height = 1013
        else:
            map_width = 2200
            map_height = 1013
        
        # Map element
        map_elem = etree.SubElement(root, "{" + CMAP_NS + "}map")
        map_elem.set("width", str(int(map_width)))
        map_elem.set("height", str(int(map_height)))
        
        # Concept list
        concept_list = etree.SubElement(map_elem, "{" + CMAP_NS + "}concept-list")
        concept_id_map = {}  # Map original IDs to generated IDs
        base_timestamp = int(datetime.now().timestamp() * 1000)
        for idx, concept in enumerate(concepts):
            # Generate a unique ID similar to CmapTools format
            concept_id = f"{base_timestamp + idx * 100}{idx}-{uuid.uuid4().hex[:8]}"
            concept_id_map[concept["id"]] = concept_id
            concept_elem = etree.SubElement(concept_list, "{" + CMAP_NS + "}concept")
            concept_elem.set("id", concept_id)
            concept_elem.set("label", concept["label"])
        
        # Linking phrase list (edge labels)
        linking_phrase_list = etree.SubElement(map_elem, "{" + CMAP_NS + "}linking-phrase-list")
        linking_phrase_map = {}  # Map edge to linking phrase ID
        
        for idx, edge in enumerate(edges):
            if isinstance(edge, dict):
                edge_label = edge.get("label", "").strip()
            else:
                edge_label = ""
            
            if edge_label:
                linking_phrase_id = f"{base_timestamp + (idx + 1000) * 100}{idx}-{uuid.uuid4().hex[:8]}"
                linking_phrase_map[idx] = linking_phrase_id
                linking_phrase_elem = etree.SubElement(
                    linking_phrase_list, "{" + CMAP_NS + "}linking-phrase"
                )
                linking_phrase_elem.set("id", linking_phrase_id)
                linking_phrase_elem.set("label", edge_label)
        
        # Connection list
        connection_list = etree.SubElement(map_elem, "{" + CMAP_NS + "}connection-list")
        
        for idx, edge in enumerate(edges):
            if isinstance(edge, dict):
                src_id = concept_id_map.get(edge.get("source"))
                dst_id = concept_id_map.get(edge.get("target"))
                edge_label = edge.get("label", "").strip()
            else:
                continue  # Skip legacy format
            
            if not src_id or not dst_id:
                continue
            
            if edge_label and idx in linking_phrase_map:
                # Connection uses a linking phrase - create two connections
                linking_phrase_id = linking_phrase_map[idx]
                
                # Connection from source concept to linking phrase
                conn1_id = f"{base_timestamp + (idx + 2000) * 100}{idx}-{uuid.uuid4().hex[:8]}"
                conn1_elem = etree.SubElement(connection_list, "{" + CMAP_NS + "}connection")
                conn1_elem.set("id", conn1_id)
                conn1_elem.set("from-id", src_id)
                conn1_elem.set("to-id", linking_phrase_id)
                
                # Connection from linking phrase to target concept
                conn2_id = f"{base_timestamp + (idx + 3000) * 100}{idx}-{uuid.uuid4().hex[:8]}"
                conn2_elem = etree.SubElement(connection_list, "{" + CMAP_NS + "}connection")
                conn2_elem.set("id", conn2_id)
                conn2_elem.set("from-id", linking_phrase_id)
                conn2_elem.set("to-id", dst_id)
            else:
                # Direct connection without linking phrase
                connection_id = f"{base_timestamp + (idx + 2000) * 100}{idx}-{uuid.uuid4().hex[:8]}"
                connection_elem = etree.SubElement(connection_list, "{" + CMAP_NS + "}connection")
                connection_elem.set("id", connection_id)
                connection_elem.set("from-id", src_id)
                connection_elem.set("to-id", dst_id)
        
        # Concept appearance list
        if layout:
            concept_appearance_list = etree.SubElement(
                map_elem, "{" + CMAP_NS + "}concept-appearance-list"
            )
            for concept in concepts:
                original_id = concept["id"]
                concept_id = concept_id_map.get(original_id)
                if concept_id and original_id in layout:
                    node = layout[original_id]
                    cx, cy = node.center_emu
                    w, h = node.size_emu
                    appearance = etree.SubElement(
                        concept_appearance_list, "{" + CMAP_NS + "}concept-appearance"
                    )
                    appearance.set("id", concept_id)
                    appearance.set("x", str(int(cx - w // 2)))
                    appearance.set("y", str(int(cy - h // 2)))
                    appearance.set("width", str(w))
                    appearance.set("height", str(h))
        
        # Linking phrase appearance list (simplified - just placeholder)
        linking_phrase_appearance_list = etree.SubElement(
            map_elem, "{" + CMAP_NS + "}linking-phrase-appearance-list"
        )
        
        # Connection appearance list (empty but required)
        connection_appearance_list = etree.SubElement(
            map_elem, "{" + CMAP_NS + "}connection-appearance-list"
        )
        
        # Style sheet list
        style_sheet_list = etree.SubElement(map_elem, "{" + CMAP_NS + "}style-sheet-list")
        style_sheet = etree.SubElement(style_sheet_list, "{" + CMAP_NS + "}style-sheet")
        style_sheet.set("id", "_Default_")
        
        map_style = etree.SubElement(style_sheet, "{" + CMAP_NS + "}map-style")
        map_style.set("background-color", "232,234,236,255")
        
        concept_style = etree.SubElement(style_sheet, "{" + CMAP_NS + "}concept-style")
        concept_style.set("font-name", "Verdana")
        concept_style.set("font-size", "12")
        concept_style.set("font-style", "bold")
        concept_style.set("font-color", "50,50,50,255")
        concept_style.set("background-color", "245,245,245,255")
        concept_style.set("border-color", "120,120,120,255")
        concept_style.set("border-thickness", "2")
        
        linking_phrase_style = etree.SubElement(
            style_sheet, "{" + CMAP_NS + "}linking-phrase-style"
        )
        linking_phrase_style.set("font-name", "Verdana")
        linking_phrase_style.set("font-size", "12")
        linking_phrase_style.set("font-color", "0,0,0,255")
        
        connection_style = etree.SubElement(style_sheet, "{" + CMAP_NS + "}connection-style")
        connection_style.set("color", "120,120,120,255")
        connection_style.set("thickness", "2")
        
        # Cmap parts list
        cmap_parts_list = etree.SubElement(root, "{" + CMAP_NS + "}cmap-parts-list")
        annotations = etree.SubElement(cmap_parts_list, "{" + CMAP_NS + "}annotations")
        etree.SubElement(annotations, "{" + CMAP_NS + "}annotation-list")
        etree.SubElement(annotations, "{" + CMAP_NS + "}annotation-appearance-list")
        
        # Convert to bytes
        xml_string = etree.tostring(
            root,
            encoding="UTF-8",
            xml_declaration=True,
            pretty_print=False,
        )
        buffer = BytesIO(xml_string)
        buffer.seek(0)
        return buffer


def concepts_from_state(state: Iterable[Tuple[str, str]]) -> List[Dict[str, str]]:
    """Helper that converts (id, label) tuples into builder-ready dicts."""
    return [{"id": concept_id, "label": label} for concept_id, label in state]
